from f_utils import u_file
from pptx import Presentation


class MyPresentationInit:

    # width of standard pptx slide
    WIDTH = 10  # 13.33
    # height of standard pptx slide
    HEIGHT = 7.5
    # 6 is the blank slide layout (without elements)
    LAYOUT_BLANK_SLIDE = 6

    def __init__(self, path: str):
        # pptx.Presentation object
        self.p = None
        # pptx.Presentation.slide object
        self.slide = None
        # Mapping {name: shape}
        self.shapes = dict()
        # paths to the pptx file
        self.path = path
        # init already existed pptx file or create list new file
        if u_file.is_exists(self.path):
            self.__init_existed_file()
        else:
            self.__init_new_file()

    def save(self, path=None) -> None:
        """
        ========================================================================
         Description: Save the PPTX-File as list initialized PPTX-Path or Save-As
                       the given new paths.
        ========================================================================
        """
        if not path:
            path = self.path
        self.p.save(path)

    def __init_existed_file(self):
        """
        ========================================================================
         Description: Init the Object in case an existing PPTX-File is given.
        ========================================================================
        """
        self.p = Presentation(self.path)
        self.slide = self.p.slides[0]

    def __init_new_file(self):
        """
        ========================================================================
         Description: Init the Object in case no existing PPTX-File is given.
        ========================================================================
        """
        self.p = Presentation()
        layout = self.p.slide_layouts[self.LAYOUT_BLANK_SLIDE]
        self.slide = self.p.slides.add_slide(layout)
