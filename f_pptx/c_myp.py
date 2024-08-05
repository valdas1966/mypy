from pptx import Presentation
from pptx.util import Inches, Pt
#from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from myp_init import MyPresentationInit
from f_utils import u_file
from f_const import u_rgb


class MyPresentation(MyPresentationInit):

    def __init__(self, path: str):
        super().__init__(path)

    def map_shapes_names(self, li: list[str]) -> None:
        self.shapes = {name: self.slide.shapes[i] for i, name in enumerate(li)}
        for i, name in enumerate(li):
            self.shapes[name] = self.slide.shapes[i]

    def organize_front(self, li):
        cur = self.slide.shapes[0]._element
        for name in li:
            shape = self.shapes[name]
            cur.addprevious(shape._element)
            cur = shape._element

    def add_text(self, params):
        left, top, width, height = params.inches
        textbox = self.slide.shapes.add_textbox(left, top, width, height)
        textbox.line.color.rgb = self.to_rgb_color(params.line_color)
        textbox.line.width = params.line_width
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = self.to_rgb_color(params.back_color)
        self.shapes[params.name] = self.slide.shapes[-1]
        tf = textbox.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  #MSO_ANCHOR.MIDDLE
        parag = tf.paragraphs[0]
        parag.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  #.CENTER
        run = parag.add_run()
        run.text = params.text
        font = run.font
        font.name = params.font
        font.size = Pt(params.text_size)
        font.color.rgb = self.to_rgb_color(params.text_color)
        font.bold = params.text_is_bold
        if params.transparency:
            shape = self.slide.shapes[-1]
            alpha = round((1-params.transparency) * 100000)
            self.__set_shape_transparency(shape, alpha)

    @classmethod
    def to_inches(cls, coor):
        left = Inches(cls.WIDTH * coor[0])
        top = Inches(cls.HEIGHT * coor[1])
        width = Inches(cls.WIDTH * coor[2])
        height = Inches(cls.HEIGHT * coor[3])
        return left, top, width, height

    @classmethod
    def to_rgb_color(cls, rgb):
        red, green, blue = rgb
        return RGBColor(red, green, blue)

    def __set_shape_transparency(self, shape, alpha):
        """ Set the transparency (alpha) of list shape"""
        ts = shape.fill._xPr.solidFill
        sF = ts.get_or_change_to_srgbClr()
        sE = self.__SubElement(sF, 'list:alpha', val=str(alpha))

    def __SubElement(self, parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    class Params:
        name = None
        inches = None
        text = str()
        font = 'Tahoma'
        text_color = u_rgb.BLACK
        text_size = 18
        text_is_bold = True
        line_width = Inches(0.01)
        line_color = u_rgb.BLACK
        back_color = u_rgb.WHITE
        transparency = 0.0
